#!/usr/bin/env python3
"""
Keynote PPT Generator — PPTX Export Engine
Converts presentation data into a polished .pptx file with Apple-quality design.

Usage:
    python3 export_pptx.py input.json output.pptx

The input JSON format:
{
  "meta": {
    "title": "Presentation Title",
    "author": "Author Name",
    "theme": "dark",           // "dark" or "light"
    "accentColor": "#6366f1"   // hex color
  },
  "slides": [
    {
      "type": "title",
      "title": "Main Title",
      "subtitle": "Subtitle text",
      "author": "Author Name"
    },
    {
      "type": "statement",
      "pill": "Key Insight",
      "title": "Bold statement text here"
    },
    {
      "type": "stats",
      "title": "Section Title",
      "stats": [
        {"number": "97%", "label": "Of Fortune 500", "sublabel": "now use AI"},
        {"number": "$4.8T", "label": "Market Size", "sublabel": "projected by 2030"},
        {"number": "10x", "label": "Productivity", "sublabel": "for early adopters"}
      ]
    },
    {
      "type": "section",
      "number": "01",
      "pill": "Chapter One",
      "title": "Section Title"
    },
    {
      "type": "content",
      "pill": "Topic",
      "title": "Content Title",
      "body": "Paragraph text...",
      "body2": "Optional second paragraph"
    },
    {
      "type": "features",
      "title": "Features Title",
      "features": [
        {"icon": "icon_text", "title": "Feature", "desc": "Description"},
      ]
    },
    {
      "type": "quote",
      "text": "Quote text here",
      "author": "Author Name",
      "role": "Title / Role"
    },
    {
      "type": "compare",
      "title": "Comparison Title",
      "left": {"title": "Before", "items": ["Item 1", "Item 2"]},
      "right": {"title": "After", "items": ["Item 1", "Item 2"]}
    },
    {
      "type": "timeline",
      "title": "Timeline Title",
      "events": [
        {"year": "2020", "title": "Event", "desc": "Description"},
      ]
    },
    {
      "type": "twocol",
      "pill": "Topic",
      "title": "Two Column Title",
      "left": {"title": "Left Title", "body": "Left content"},
      "right": {"title": "Right Title", "body": "Right content"}
    },
    {
      "type": "closing",
      "title": "Closing Title",
      "subtitle": "Closing subtitle",
      "cta": "Call to action text",
      "contact": "email@example.com"
    }
  ]
}
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Design Tokens ──────────────────────────────────────────

THEMES = {
    "dark": {
        "bg": RGBColor(0x0A, 0x0A, 0x1A),
        "bg_card": RGBColor(0x16, 0x16, 0x3A),
        "text": RGBColor(0xF1, 0xF5, 0xF9),
        "text_secondary": RGBColor(0x94, 0xA3, 0xB8),
        "text_muted": RGBColor(0x64, 0x74, 0x8B),
        "accent": RGBColor(0x63, 0x66, 0xF1),
        "accent2": RGBColor(0x8B, 0x5C, 0xF6),
        "accent_cyan": RGBColor(0x22, 0xD3, 0xEE),
        "accent_rose": RGBColor(0xF4, 0x3F, 0x5E),
        "accent_emerald": RGBColor(0x10, 0xB9, 0x81),
        "accent_amber": RGBColor(0xF5, 0x9E, 0x0B),
        "card_border": RGBColor(0x2D, 0x2D, 0x5A),
    },
    "light": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_card": RGBColor(0xF8, 0xFA, 0xFC),
        "text": RGBColor(0x0F, 0x17, 0x2A),
        "text_secondary": RGBColor(0x64, 0x74, 0x8B),
        "text_muted": RGBColor(0x94, 0xA3, 0xB8),
        "accent": RGBColor(0x63, 0x66, 0xF1),
        "accent2": RGBColor(0x8B, 0x5C, 0xF6),
        "accent_cyan": RGBColor(0x06, 0xB6, 0xD4),
        "accent_rose": RGBColor(0xE1, 0x1D, 0x48),
        "accent_emerald": RGBColor(0x05, 0x9F, 0x73),
        "accent_amber": RGBColor(0xD9, 0x77, 0x06),
        "card_border": RGBColor(0xE2, 0xE8, 0xF0),
    },
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Inter"
FALLBACK_FONT = "Helvetica Neue"


def hex_to_rgb(hex_str):
    """Convert #RRGGBB to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=None, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_NAME, line_spacing=1.4):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, radius=Inches(0.2)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Set corner radius
    shape.adjustments[0] = 0.08
    return shape


def add_pill(slide, left, top, text, theme, color_key="accent"):
    """Add a pill/tag shape."""
    w, h = Inches(2.2), Inches(0.4)
    pill = add_rounded_rect(slide, left, top, w, h, theme["bg_card"], theme.get(color_key, theme["accent"]))
    pill.adjustments[0] = 0.5  # Full round
    # Text on pill
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.name = FONT_NAME
    p.font.color.rgb = theme.get(color_key, theme["accent"])
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return pill


def add_divider(slide, left, top, width=Inches(1), color=None):
    """Add a gradient-like divider line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or RGBColor(0x63, 0x66, 0xF1)
    shape.line.fill.background()
    return shape


# ── Slide Builders ─────────────────────────────────────────

def build_title_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, theme["bg"])

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(2.0),
                 data.get("title", ""), font_size=54, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    if data.get("subtitle"):
        add_text_box(slide, Inches(2.5), Inches(4.2), Inches(8), Inches(1.0),
                     data["subtitle"], font_size=22, color=theme["text_secondary"],
                     alignment=PP_ALIGN.CENTER)

    # Author
    if data.get("author"):
        add_text_box(slide, Inches(2.5), Inches(5.5), Inches(8), Inches(0.5),
                     data["author"].upper(), font_size=12, color=theme["text_muted"],
                     bold=True, alignment=PP_ALIGN.CENTER)

    # Accent bar at top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.5), Inches(4), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()


def build_statement_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    if data.get("pill"):
        add_pill(slide, Inches(4.5), Inches(2.2), data["pill"], theme, "accent_cyan")

    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(2.5),
                 data.get("title", ""), font_size=42, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.2)


def build_stats_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
                 data.get("title", ""), font_size=32, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    stats = data.get("stats", [])
    n = len(stats)
    if n == 0:
        return

    card_w = Inches(3.5)
    card_h = Inches(3.0)
    gap = Inches(0.4)
    total_w = n * card_w + (n - 1) * gap
    start_x = (SLIDE_WIDTH - total_w) / 2

    colors = [theme["accent"], theme["accent_cyan"], theme["accent_rose"],
              theme["accent_emerald"], theme["accent_amber"]]

    for i, stat in enumerate(stats):
        x = start_x + i * (card_w + gap)
        y = Inches(2.4)
        card = add_rounded_rect(slide, x, y, card_w, card_h, theme["bg_card"], theme["card_border"])

        add_text_box(slide, x, y + Inches(0.5), card_w, Inches(1.2),
                     str(stat.get("number", "")), font_size=52,
                     color=colors[i % len(colors)], bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, y + Inches(1.6), card_w, Inches(0.4),
                     stat.get("label", "").upper(), font_size=11,
                     color=theme["text_secondary"], bold=True, alignment=PP_ALIGN.CENTER)

        if stat.get("sublabel"):
            add_text_box(slide, x, y + Inches(2.1), card_w, Inches(0.4),
                         stat["sublabel"], font_size=10,
                         color=theme["text_muted"], alignment=PP_ALIGN.CENTER)


def build_section_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # Big faded number
    if data.get("number"):
        add_text_box(slide, Inches(3), Inches(0.5), Inches(7), Inches(3),
                     str(data["number"]), font_size=120,
                     color=theme["text_muted"], bold=True, alignment=PP_ALIGN.CENTER)

    if data.get("pill"):
        add_pill(slide, Inches(5), Inches(2.8), data["pill"], theme, "accent_emerald")

    add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(2.0),
                 data.get("title", ""), font_size=52, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)


def build_content_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # Left column - text
    x = Inches(1)
    if data.get("pill"):
        add_pill(slide, x, Inches(1.2), data["pill"], theme)

    add_text_box(slide, x, Inches(1.8), Inches(5.5), Inches(1.2),
                 data.get("title", ""), font_size=36, color=theme["text"], bold=True)

    add_divider(slide, x, Inches(3.1), Inches(1))

    if data.get("body"):
        add_text_box(slide, x, Inches(3.5), Inches(5.5), Inches(2.0),
                     data["body"], font_size=16, color=theme["text_secondary"],
                     line_spacing=1.6)

    if data.get("body2"):
        add_text_box(slide, x, Inches(5.0), Inches(5.5), Inches(1.5),
                     data["body2"], font_size=16, color=theme["text_secondary"],
                     line_spacing=1.6)

    # Right column - visual card
    rx = Inches(7.5)
    card = add_rounded_rect(slide, rx, Inches(1.8), Inches(5), Inches(4.5),
                            theme["bg_card"], theme["card_border"])


def build_features_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
                 data.get("title", ""), font_size=32, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    features = data.get("features", [])
    cols = 3
    rows = (len(features) + cols - 1) // cols
    card_w = Inches(3.6)
    card_h = Inches(2.4)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    total_w = cols * card_w + (cols - 1) * gap_x
    start_x = (SLIDE_WIDTH - total_w) / 2
    start_y = Inches(2.0)

    colors = [theme["accent"], theme["accent_cyan"], theme["accent2"],
              theme["accent_emerald"], theme["accent_rose"], theme["accent_amber"]]

    for i, feat in enumerate(features):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rounded_rect(slide, x, y, card_w, card_h, theme["bg_card"], theme["card_border"])

        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), Inches(0.4),
                     feat.get("icon", ""), font_size=24, color=colors[i % len(colors)])

        add_text_box(slide, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), Inches(0.4),
                     feat.get("title", ""), font_size=16, color=theme["text"], bold=True)

        add_text_box(slide, x + Inches(0.3), y + Inches(1.3), card_w - Inches(0.6), Inches(1.0),
                     feat.get("desc", ""), font_size=12, color=theme["text_secondary"],
                     line_spacing=1.5)


def build_quote_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # Big quote mark
    add_text_box(slide, Inches(4.5), Inches(0.8), Inches(4), Inches(2),
                 "\u201C", font_size=160, color=theme["text_muted"],
                 alignment=PP_ALIGN.CENTER)

    # Quote text
    add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(2.5),
                 data.get("text", ""), font_size=28, color=theme["text"],
                 italic=True, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Attribution
    author_text = data.get("author", "")
    if data.get("role"):
        author_text += f" \u2014 {data['role']}"
    add_text_box(slide, Inches(2), Inches(5.2), Inches(9), Inches(0.5),
                 author_text, font_size=16, color=theme["text_secondary"],
                 alignment=PP_ALIGN.CENTER)


def build_compare_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
                 data.get("title", ""), font_size=32, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    card_w = Inches(5.2)
    card_h = Inches(4.5)
    left_x = Inches(0.8)
    right_x = Inches(7.2)
    y = Inches(2.0)

    # Left card
    add_rounded_rect(slide, left_x, y, card_w, card_h, theme["bg_card"], theme["card_border"])
    left = data.get("left", {})
    add_text_box(slide, left_x + Inches(0.4), y + Inches(0.3), card_w - Inches(0.8), Inches(0.5),
                 left.get("title", ""), font_size=22, color=theme["accent_rose"], bold=True)
    for j, item in enumerate(left.get("items", [])):
        add_text_box(slide, left_x + Inches(0.6), y + Inches(1.0 + j * 0.6), card_w - Inches(1), Inches(0.5),
                     f"\u2022  {item}", font_size=14, color=theme["text_secondary"])

    # VS
    add_text_box(slide, Inches(6.0), Inches(3.5), Inches(1.2), Inches(0.5),
                 "\u2192", font_size=28, color=theme["text_muted"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Right card
    add_rounded_rect(slide, right_x, y, card_w, card_h, theme["bg_card"], theme["card_border"])
    right = data.get("right", {})
    add_text_box(slide, right_x + Inches(0.4), y + Inches(0.3), card_w - Inches(0.8), Inches(0.5),
                 right.get("title", ""), font_size=22, color=theme["accent_emerald"], bold=True)
    for j, item in enumerate(right.get("items", [])):
        add_text_box(slide, right_x + Inches(0.6), y + Inches(1.0 + j * 0.6), card_w - Inches(1), Inches(0.5),
                     f"\u2022  {item}", font_size=14, color=theme["text_secondary"])


def build_timeline_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
                 data.get("title", ""), font_size=32, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    events = data.get("events", [])
    n = len(events)
    if n == 0:
        return

    # Horizontal line
    line_y = Inches(3.2)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.5), line_y, Inches(10), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = theme["card_border"]
    line.line.fill.background()

    item_w = Inches(10) / n
    start_x = Inches(1.5)

    for i, evt in enumerate(events):
        cx = start_x + i * item_w + item_w / 2

        # Dot
        dot_size = Inches(0.18)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx - dot_size / 2, line_y - dot_size / 2 + Pt(1),
                                     dot_size, dot_size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme["accent"]
        dot.line.fill.background()

        # Year
        add_text_box(slide, cx - Inches(1), Inches(3.6), Inches(2), Inches(0.3),
                     evt.get("year", ""), font_size=11, color=theme["accent"],
                     bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, cx - Inches(1), Inches(3.9), Inches(2), Inches(0.4),
                     evt.get("title", ""), font_size=15, color=theme["text"],
                     bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, cx - Inches(1.2), Inches(4.4), Inches(2.4), Inches(1.2),
                     evt.get("desc", ""), font_size=11, color=theme["text_secondary"],
                     alignment=PP_ALIGN.CENTER, line_spacing=1.4)


def build_twocol_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    if data.get("pill"):
        add_pill(slide, Inches(1), Inches(0.8), data["pill"], theme, "accent_rose")

    add_text_box(slide, Inches(1), Inches(1.4), Inches(11), Inches(0.8),
                 data.get("title", ""), font_size=32, color=theme["text"], bold=True)

    card_w = Inches(5.5)
    card_h = Inches(3.8)
    y = Inches(2.8)

    # Left card
    left = data.get("left", {})
    add_rounded_rect(slide, Inches(0.8), y, card_w, card_h, theme["bg_card"], theme["card_border"])
    add_text_box(slide, Inches(1.2), y + Inches(0.4), card_w - Inches(0.8), Inches(0.5),
                 left.get("title", ""), font_size=18, color=theme["accent_cyan"], bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(1.1), card_w - Inches(0.8), Inches(2.2),
                 left.get("body", ""), font_size=14, color=theme["text_secondary"],
                 line_spacing=1.6)

    # Right card
    right = data.get("right", {})
    add_rounded_rect(slide, Inches(6.8), y, card_w, card_h, theme["bg_card"], theme["card_border"])
    add_text_box(slide, Inches(7.2), y + Inches(0.4), card_w - Inches(0.8), Inches(0.5),
                 right.get("title", ""), font_size=18, color=theme["accent_emerald"], bold=True)
    add_text_box(slide, Inches(7.2), y + Inches(1.1), card_w - Inches(0.8), Inches(2.2),
                 right.get("body", ""), font_size=14, color=theme["text_secondary"],
                 line_spacing=1.6)


def build_closing_slide(prs, data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 data.get("title", ""), font_size=48, color=theme["text"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    if data.get("subtitle"):
        add_text_box(slide, Inches(2.5), Inches(3.8), Inches(8), Inches(0.8),
                     data["subtitle"], font_size=18, color=theme["text_secondary"],
                     alignment=PP_ALIGN.CENTER)

    if data.get("cta"):
        # CTA button-like shape
        btn_w = Inches(4)
        btn_h = Inches(0.65)
        btn = add_rounded_rect(slide,
                               (SLIDE_WIDTH - btn_w) / 2, Inches(4.8),
                               btn_w, btn_h, theme["accent"])
        btn.adjustments[0] = 0.5
        tf = btn.text_frame
        p = tf.paragraphs[0]
        p.text = data["cta"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = FONT_NAME
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if data.get("contact"):
        add_text_box(slide, Inches(2.5), Inches(5.8), Inches(8), Inches(0.5),
                     data["contact"], font_size=12, color=theme["text_muted"],
                     alignment=PP_ALIGN.CENTER)


# ── Main Builder ───────────────────────────────────────────

BUILDERS = {
    "title": build_title_slide,
    "statement": build_statement_slide,
    "stats": build_stats_slide,
    "section": build_section_slide,
    "content": build_content_slide,
    "features": build_features_slide,
    "quote": build_quote_slide,
    "compare": build_compare_slide,
    "timeline": build_timeline_slide,
    "twocol": build_twocol_slide,
    "closing": build_closing_slide,
}


def build_presentation(data):
    """Build a complete .pptx presentation from JSON data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    meta = data.get("meta", {})
    theme_name = meta.get("theme", "dark")
    theme = THEMES.get(theme_name, THEMES["dark"])

    # Override accent color if provided
    if meta.get("accentColor"):
        theme["accent"] = hex_to_rgb(meta["accentColor"])

    for slide_data in data.get("slides", []):
        slide_type = slide_data.get("type", "content")
        builder = BUILDERS.get(slide_type, build_content_slide)
        builder(prs, slide_data, theme)

    return prs


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 export_pptx.py input.json output.pptx")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    with open(input_path, "r") as f:
        data = json.load(f)

    prs = build_presentation(data)
    prs.save(output_path)
    print(f"Saved presentation to: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Theme: {data.get('meta', {}).get('theme', 'dark')}")


if __name__ == "__main__":
    main()
